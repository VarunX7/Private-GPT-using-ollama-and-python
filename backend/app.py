from flask import Flask, request, jsonify
from flask_cors import CORS
from langchain_community.document_loaders import UnstructuredPDFLoader
from langchain_ollama import OllamaEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain.prompts import ChatPromptTemplate, PromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_ollama import ChatOllama
from langchain_core.runnables import RunnablePassthrough
from langchain.retrievers.multi_query import MultiQueryRetriever
from langchain.schema import Document
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd
import hashlib
import ollama
import os

app = Flask(__name__)
CORS(app)

# stimulate db for storing metadata
processed_documents = {}

# Preload embeddings model to save time
ollama.pull("nomic-embed-text")

# database for storing chunks and embeddings
vector_db = Chroma

# function to calculate hash of the file provided ........................................
def calculate_hash(file_path):
    hasher = hashlib.sha256()
    with open(file_path, "rb") as file:
        while chunk := file.read(8192):
            hasher.update(chunk)
    return



# function to load content based on file type ..............................................
def load_file_content(file_path):

    ext = os.path.splitext(file_path)[1].lower()
    content = ""

    if ext == ".pdf":
        loader = UnstructuredPDFLoader(file_path=file_path)
        content=loader.load()
    elif ext == ".docx":
        doc = DocxDocument(file_path)
        combined_content = "\n".join([para.text.strip() for para in doc.paragraphs if para.text.strip()])
        content = [Document(page_content=combined_content)]
    elif ext == ".pptx":
        ppt = Presentation(file_path)
        combined_content = "\n".join(
            [slide.notes_slide.notes_text_frame.text.strip() for slide in ppt.slides if slide.notes_slide and slide.notes_slide.notes_text_frame]
        )
        content = [Document(page_content=combined_content)]
    elif ext in [".csv", ".xlsx"]:
        df = pd.read_csv(file_path) if ext == ".csv" else pd.read_excel(file_path)
        combined_content = "\n".join([" | ".join([str(cell) for cell in row]) for _, row in df.iterrows()])
        content = [Document(page_content=combined_content)]
    else:
        raise ValueError(f"Unsupported file type: {ext}")
    
    return content


# Function to process and store file into the vector database ............................
def process_and_store_doc(file_path, doc_hash):

    global vector_db

    #Load the PDF content
    data = load_file_content(file_path)
    print("Done loading !")

    # Split the content into chunks
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1200, chunk_overlap=300 )
    chunks = text_splitter.split_documents(data)
    print("done splitting !")

    # Add chuks to the vector database
    vector_db = Chroma.from_documents(
        documents=chunks,
        embedding=OllamaEmbeddings(model="nomic-embed-text"),
        collection_name="simple_rag",
    )   

    # Store metadata 
    processed_documents[os.path.basename(file_path)] = {
        "hash":doc_hash,
        "vector_db": vector_db
    }


# Function to skip the processing if the file has already been processed ...................
def does_file_exist(pdf, temp_pdf_path):

    global vector_db 

    doc_hash = calculate_hash(temp_pdf_path)

    if pdf.filename in processed_documents:
        if processed_documents[pdf.filename]['hash'] == doc_hash:
            print("Using previously stored chunks...")
            vector_db = processed_documents[pdf.filename]['vector_db']            
        else:
            print("Document modified, Reprocessing...")
            process_and_store_doc(temp_pdf_path, doc_hash)
    else:
        print("New document. Processing and storing...")
        process_and_store_doc(temp_pdf_path, doc_hash)

# ----------------------------------- FILE UPLOAD ROUTE ---------------------------------------

@app.route('/upload', methods=['POST'])
def upload():
    pdf = request.files['file']

    # Save the PDF to a temporary location...
    temp_pdf_path = os.path.join("temp", pdf.filename)
    os.makedirs("temp", exist_ok = True)
    pdf.save(temp_pdf_path)

    #Calculate hash for document...
    does_file_exist(pdf, temp_pdf_path)

    return jsonify({"message": "File uploaded and text extracted successfully :)"})

# -------------------------------- ANSWERING QUESTION ROUTE -----------------------------------

@app.route('/ask-question', methods=['POST'])
def ask_question():
    global vector_db

    print("running ask_question")
    try:
        if 'question' not in request.form:
            return jsonify({"error": "Ask a question"}), 400
        
        question = request.form['question']
        
        # Set up retrieval and llm...
        model = "llama3.2"
        llm = ChatOllama(model=model)

        QUERY_PROMPT = PromptTemplate(
            input_variables=["question"],
            template="""You are an AI language model assistant. Your task is to generate five different versions of the user's question to retrieve relevant documents from the database. By generating multiple perpectives on the user question, your goal is to help the user overcome some of the limitations of the distance-based similarity search. Provide these alternative questions seperated by new lines. 
            Original Question: {question}"""
        )

        retriever = MultiQueryRetriever.from_llm(
            vector_db.as_retriever(), llm, prompt=QUERY_PROMPT
        )

        # RAG prompt
        template = """Answer the questions based on only the following context: 
        {context}
        Question: {question}
        """

        prompt = ChatPromptTemplate.from_template(template)

        chain = (
            {"context": retriever, "question": RunnablePassthrough()}
            |prompt
            |llm
            |StrOutputParser()
        )

        # Get the result
        res = chain.invoke(input=(question,))
        print(res)

        # Clean up temporary files
        # os.remove(temp_pdf_path)

        # Return the answer
        return jsonify({"answer": res})
    
    except Exception as e:
        print(str(e))
        return jsonify({"error": str(e)}), 500
    
# ------------------------------------------- MAIN ---------------------------------------------
    
if __name__ == '__main__':
    app.run(debug=True)
